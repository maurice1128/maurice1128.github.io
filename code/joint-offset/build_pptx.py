"""Lab-presentation deck: figure-heavy, text-light.
Palette matches the paper's figures so slides and plots read as one system.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

INK   = RGBColor(0x1A, 0x1A, 0x1A)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x24, 0x28, 0x2B)
HARM  = RGBColor(0xB2, 0x3A, 0x48)
HELP  = RGBColor(0x2E, 0x6E, 0x8E)
MUTE  = RGBColor(0x6E, 0x6E, 0x6E)          # 4.6:1 on white (was 8A8A8A, 3.5:1)
MUTE_D = RGBColor(0xB4, 0xBA, 0xBE)         # for dark backgrounds
FAINT = RGBColor(0xD8, 0xD8, 0xD8)

SW, SH = 13.333, 7.5
FIG = "D:/ROWV_paper/figures/"
TITLE_X, TITLE_Y = 0.9, 0.72                # every slide title, figure slides included

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
BLANK = prs.slide_layouts[6]


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def text(s, txt, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         font="Calibri", italic=False, space=6):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = font
    return tb


def title(s, txt, size=32, color=INK):
    return text(s, txt, TITLE_X, TITLE_Y, 11.8, 0.9, size=size, bold=True,
                color=color, font="Cambria")


def picture(s, fn, top, bottom, maxw=12.2):
    iw, ih = Image.open(FIG + fn).size
    ar = iw / ih
    w = min(maxw, (bottom - top) * ar)
    h = w / ar
    return s.shapes.add_picture(FIG + fn, Inches((SW - w) / 2),
                                Inches(top + (bottom - top - h) / 2),
                                width=Inches(w), height=Inches(h))


def stat(s, value, label, x, y, w, color=INK, vsize=72):
    text(s, value, x, y, w, 1.3, size=vsize, bold=True, color=color,
         align=PP_ALIGN.CENTER, font="Cambria")
    text(s, label, x, y + 1.15, w, 0.9, size=13, color=MUTE, align=PP_ALIGN.CENTER)


def card(s, x, y, w, h, col, lw=2):
    c = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    c.line.color.rgb = col; c.line.width = Pt(lw)
    c.shadow.inherit = False
    c.text_frame.text = ""
    return c


DOT, DASH, MINUS, DEG = "\u00b7", "\u2014", "\u2212", "\u00b0"

# --- counts derived from the manuscript, never hardcoded (drift guard) ---
import io, re, sys; sys.path.insert(0, "D:/ROWV_paper")
from paper_stats import word_count as _wc, reference_count as _rc, title as _title
_M_INTER = int(re.search(r"Family size m = (\d+)",
                        io.open("D:/BioCV/BIOCV_INTERACTION_FDR.txt",
                                encoding="utf-8").read()).group(1))
_WORDS, _REFS = _wc(), _rc()

# ---------------------------------------------------------------- 1 title
s = slide(DARK)
# Wrapped from the manuscript title, never retyped -- it lived in three files at once and drifted.
_t = _title()
_mid = _t.rfind(" ", 0, len(_t) // 2 + 12)
text(s, _t[:_mid] + "\n" + _t[_mid + 1:],
     0.9, 2.15, 11.5, 2.4, size=40, bold=True, color=PAPER, font="Cambria", space=10)
text(s, f"Multi-view markerless motion capture   {DOT}   BioCV dataset   {DOT}   "
        "11 participants, 104 walking trials",
     0.9, 4.8, 11.5, 0.5, size=15, color=MUTE_D)
text(s, f"Mu-Hua Wang   {DOT}   advisor: Hui-Ting Shih   {DOT}   NYCU",
     0.9, 5.5, 11.5, 0.4, size=13, color=RGBColor(0x8E, 0x96, 0x9B))

# ---------------------------------------------------------------- 2 setup
s = slide()
title(s, "How markerless systems get validated", size=34)
text(s, f"Tuned and reported against 3D joint-centre position error {DASH} "
        "inherited from computer vision.", 0.9, 1.72, 11.5, 0.5, size=17, color=MUTE)
for i, (word, lab, col) in enumerate([
        ("millimetres", "what the pipeline is tuned on\n3D position error, averaged over joints", HELP),
        ("degrees", "what the gait lab actually reports\njoint angle", HARM)]):
    x = 1.2 + i * 5.5
    card(s, x, 2.55, 4.9, 2.5, col, lw=2)
    text(s, word, x, 2.95, 4.9, 1.0, size=54, bold=True, color=col,
         align=PP_ALIGN.CENTER, font="Cambria")
    text(s, lab, x + 0.3, 4.1, 4.3, 1.0, size=15, color=INK,
         align=PP_ALIGN.CENTER, space=3)
text(s, "Same pipeline, same trials. Two ways of scoring the very same output.",
     0.9, 5.4, 11.5, 0.5, size=15, color=MUTE, align=PP_ALIGN.CENTER)
text(s, "Are they interchangeable?", 0.9, 6.15, 11.5, 0.6, size=26, bold=True,
     align=PP_ALIGN.CENTER, font="Cambria")

# ---------------------------------------------------------------- 3 what we did
s = slide()
title(s, "One pipeline, five decisions, varied one at a time")
items = [("confidence weighting", "trust each camera by its detector score", HELP),
         ("inverse-distance weighting", "trust nearer cameras more", HELP),
         (f"outlier rejection {DASH} two criteria", "DISCARD cameras whose residual looks wrong", HARM),
         ("Gauss\u2013Newton refinement", "iterative geometric fit", MUTE),
         ("joint-centre offset correction", "subtract a published offset table", MUTE)]
for i, (lab, sub, col) in enumerate(items):
    x = 0.9 + (i % 3) * 4.0
    y = 1.95 + (i // 3) * 1.55
    card(s, x, y, 3.6, 1.3, col, lw=1.5)
    text(s, lab, x + 0.2, y + 0.18, 3.2, 0.4, size=14, bold=True, color=INK,
         align=PP_ALIGN.CENTER)
    text(s, sub, x + 0.2, y + 0.62, 3.2, 0.6, size=11, color=MUTE,
         align=PP_ALIGN.CENTER, space=2)
text(s, "Same detections.\nSame calibration.\nSame trials.\n\n"
        f"Scored on 3D position\nand on three joint angles {DASH}\nfrom the identical output.",
     8.95, 3.55, 3.5, 1.8, size=14, color=INK, space=1)
text(s, f"50 contrasts  {DOT}  exact sign-flip permutation over the 11 participants  "
        f"{DOT}  Benjamini\u2013Hochberg", 0.9, 5.75, 11.5, 0.5, size=15, color=MUTE)
text(s, f"The interaction tests form their own family of {_M_INTER}, corrected separately.",
     0.9, 6.25, 11.5, 0.4, size=13, color=MUTE)

# ---------------------------------------------------------------- 4 ground truth
s = slide()
title(s, "What we compare against")
text(s, "15-camera Qualisys, 200 Hz, custom 38-marker 6-DoF set. Same trials, hardware-synchronised.",
     0.9, 1.72, 11.5, 0.5, size=17, color=MUTE)
for i, (val, lab, col) in enumerate([
        ("3.6 mm", "the reference's own\ncross-system alignment error", MUTE),
        ("3 to 6 mm", "how far the joint-centre offset\nswings with posture", MUTE),
        ("0.5 mm", "the effects we call\nsignificant", HARM)]):
    x = 1.0 + i * 4.0
    text(s, val, x, 2.5, 3.5, 1.0, size=44, bold=True, color=col,
         align=PP_ALIGN.CENTER, font="Cambria")
    text(s, lab, x, 3.6, 3.5, 0.9, size=14, color=MUTE, align=PP_ALIGN.CENTER, space=2)
box = card(s, 1.0, 4.85, 11.3, 1.75, HELP, lw=2)
text(s, "So how can a 0.5 mm effect mean anything?", 1.4, 5.1, 10.5, 0.5,
     size=20, bold=True, color=HELP, font="Cambria")
text(s, "Every contrast is a PAIRED difference on the same frames, same detections, same cameras.\n"
        "Whatever the reference gets wrong, both arms inherit equally \u2014 it cancels in the difference.",
     1.4, 5.65, 10.5, 0.9, size=15, color=INK, space=3)

# ---------------------------------------------------------------- 5 the finding
s = slide(DARK)
text(s, "The finding", 0.9, 2.5, 11.5, 0.6, size=22, color=MUTE_D, align=PP_ALIGN.CENTER)
text(s, "A null in the average was two\nsignificant effects cancelling",
     0.9, 3.15, 11.5, 2.0, size=42, bold=True, color=PAPER,
     align=PP_ALIGN.CENTER, font="Cambria", space=10)

# ---------------------------------------------------------------- 6 fig 2
s = slide()
title(s, "Same intervention, opposite effects at different joints", size=28)
picture(s, "fig2_perjoint.png", 1.75, 6.35, maxw=11.2)
text(s, f"Object-space rejection at matched retention  {DOT}  3D position error, millimetres  "
        f"{DOT}  positive = that criterion is better",
     0.9, 6.55, 11.8, 0.5, size=14, color=MUTE, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- 7 the numbers
s = slide()
title(s, "What the average hides", size=34)
text(s, "3D position error in millimetres. Positive = the object-space criterion is better.",
     0.9, 1.72, 11.5, 0.5, size=16, color=MUTE)
# q, not p: the per-joint tests are BH-corrected within their own family of twelve and the
# manuscript quotes q. The deck kept the unadjusted p for a round after that change.
stat(s, "+0.53", "HIP\nsignificantly better\nq = 0.004", 1.0, 2.55, 3.4, HELP)
stat(s, f"{MINUS}0.66", "KNEE\nsignificantly worse\nq = 0.004", 5.0, 2.55, 3.4, HARM)
stat(s, "+0.03", "THREE-JOINT MEAN\nnot significant\nq = 0.65", 9.0, 2.55, 3.4, MUTE)
text(s, "Report the mean, and you conclude the setting does nothing.",
     0.9, 5.85, 11.5, 0.6, size=24, bold=True, align=PP_ALIGN.CENTER, font="Cambria")

# ---------------------------------------------------------------- 8 sign reversal
s = slide()
title(s, "And one contrast reverses outright")
text(s, "Inverse-distance weighting, same trials, scored on two different joint sets",
     0.9, 1.68, 11.5, 0.5, size=17, color=MUTE)
for i, (val, lab, col, sig) in enumerate([
        ("+0.256 mm", "on ankles + wrists\nsignificantly BETTER", HELP, "q = 0.038"),
        (f"{MINUS}0.691 mm", "on hip, knee, ankle\nsignificantly WORSE", HARM, "q = 0.004")]):
    x = 1.9 + i * 5.2
    card(s, x, 2.5, 4.3, 2.85, col)
    text(s, val, x, 2.82, 4.3, 1.1, size=48, bold=True, color=col,
         align=PP_ALIGN.CENTER, font="Cambria")
    text(s, lab, x, 4.02, 4.3, 0.9, size=15, color=INK, align=PP_ALIGN.CENTER, space=2)
    text(s, sig, x, 4.85, 4.3, 0.4, size=13, color=MUTE, align=PP_ALIGN.CENTER)
text(s, f"Two opposite significant conclusions from one dataset {DASH}\n"
        "the joint set decides which one you report.",
     0.9, 5.75, 11.5, 1.0, size=20, bold=True, align=PP_ALIGN.CENTER,
     font="Cambria", space=4)

# ---------------------------------------------------------------- 9 why it matters
s = slide()
title(s, "Why it matters", size=34)
rows = ["The knee is the joint that criterion degrades.",
        "Knee flexion is the angle at the knee.",
        f"So the angle degrades too {DASH} 0.2 to 0.4{DEG} for every camera you discard."]
for i, t in enumerate(rows):
    c = HARM if i == 2 else INK
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(2.35 + i * 1.05),
                             Inches(0.44), Inches(0.44))
    dot.fill.solid(); dot.fill.fore_color.rgb = c
    dot.line.fill.background(); dot.shadow.inherit = False
    p = dot.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i + 1)
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = PAPER; r.font.name = "Calibri"
    text(s, t, 2.35, 2.4 + i * 1.05, 9.8, 0.7, size=21, color=c, bold=(i == 2))
card(s, 1.5, 5.6, 10.6, 1.15, HARM, lw=2)
text(s, "Judge by the pooled number and you keep a setting\nthat is harming the outcome you report.",
     1.9, 5.82, 9.8, 0.9, size=20, bold=True, font="Cambria", color=HARM, space=2)

# ---------------------------------------------------------------- 10 fig 1
s = slide()
title(s, "Reweighting reaches the angle. Discarding does not.", size=28)
picture(s, "fig1_dissociation.png", 1.65, 7.25, maxw=11.6)

# ---------------------------------------------------------------- 11 offset backfires
s = slide()
title(s, "The obvious fix makes the angle worse")
text(s, f"Publish a joint-centre offset correction table {DASH} "
        "the remedy everyone reaches for", 0.9, 1.68, 11.5, 0.5, size=17, color=MUTE)
stat(s, "+9.85 mm", "3D POSITION\nlargest gain in the study", 1.9, 2.75, 4.2, HELP, vsize=50)
ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.45), Inches(3.15),
                        Inches(0.95), Inches(0.5))
ar.fill.solid(); ar.fill.fore_color.rgb = FAINT
ar.line.fill.background(); ar.shadow.inherit = False
stat(s, f"{MINUS}0.61{DEG}", "KNEE FLEXION\nsignificantly worse", 7.7, 2.75, 4.2, HARM, vsize=50)
text(s, f"The offset is not a constant {DASH} it swings 3 to 6 mm with posture across the gait cycle.\n"
        "A static table subtracts a fixed number from a moving quantity.",
     1.3, 5.5, 11.0, 1.1, size=17, align=PP_ALIGN.CENTER, space=4)

# ---------------------------------------------------------------- 12 posture
s = slide()
title(s, "The offset moves with posture", size=28)
picture(s, "fig5_posture.png", 1.65, 6.3, maxw=8.2)
# q, not p: the posture family is BH-corrected within its own thirteen tests, and the deck
# kept the unadjusted values for a round after the paper moved to q.
text(s, f"hip +6.19 mm (q = 0.007)   {DOT}   knee {MINUS}2.82 mm (q = 0.029)   {DOT}   "
        f"ankle {MINUS}4.27 mm (q = 0.013)",
     0.9, 6.5, 11.8, 0.5, size=14, color=MUTE, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- 13 takeaways
s = slide()
title(s, "Two things to change", size=34)
text(s, "Neither depends on effect size. A sign reversal is a sign reversal at any magnitude.",
     0.9, 1.72, 11.5, 0.5, size=17, color=MUTE)
for i, (head, body, col) in enumerate([
        ("Measure position on the joints\nyour reported outcome uses",
         "A gain measured elsewhere may not exist there,\nand a null there may be two opposed effects.",
         HELP),
        ("Test the dissociation,\nnot two significance verdicts",
         "Significant in one outcome and not the other\nis not evidence that they differ.", HARM)]):
    x = 0.9 + i * 6.1
    card(s, x, 2.5, 5.6, 3.3, col)
    text(s, head, x + 0.4, 2.88, 4.8, 1.5, size=22, bold=True, color=col,
         font="Cambria", space=4)
    text(s, body, x + 0.4, 4.5, 4.8, 1.2, size=15, color=INK, space=3)
text(s, "Effects here are 0.004 to 0.61 degrees, below any clinical threshold. "
        "The claim is about how validation is reported,\nnot about whether these settings change a diagnosis.",
     0.9, 6.15, 11.5, 0.9, size=15, color=MUTE, align=PP_ALIGN.CENTER, space=3)

# ---------------------------------------------------------------- 14 status
s = slide(DARK)
text(s, "Where it stands", 0.9, 2.15, 11.5, 0.7, size=26, color=PAPER,
     bold=True, font="Cambria")
for i, (k, v) in enumerate([
        ("Target", "Journal of Biomechanics \u2014 Full Paper"),
        ("Draft", f"{_WORDS:,} words  {DOT}  3 tables  {DOT}  4 figures  {DOT}  {_REFS} references"),
        ("Open", f"funding  {DOT}  CRediT  {DOT}  code deposit")]):
    y = 3.25 + i * 0.85
    text(s, k, 0.9, y, 2.0, 0.5, size=16, color=MUTE_D, bold=True)
    text(s, v, 3.1, y, 9.4, 0.6, size=19, color=PAPER)
text(s, "One rig, one detector, 11 participants.",
     0.9, 6.15, 11.5, 0.5, size=15, color=MUTE_D, italic=True)

prs.save("D:/ROWV_paper/ROWV_lab_talk.pptx")
print("wrote ROWV_lab_talk.pptx —", len(prs.slides._sldIdLst), "slides")
